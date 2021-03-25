from pptx import Presentation
import re

def read_ppt(ppt_file_path):
    text_runs = []
    prs = Presentation(ppt_file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
    #특수문자 치환
    text_runs = [re.sub('[-−+=_:;#/?:$/|\@^*[\]{\}(\)<\`\>~0123456789]', ' ', i) for i in text_runs]
    #공백항 제거, 단어별로 분리
    text_runs = ' '.join(text_runs).split()
    text_runs = ' '.join(text_runs)
    return text_runs


# prs = Presentation('C:/Users/이윤정/Desktop/캡디/참고파일/3조 최종데모 발표 자료.pptx')

# f = open("C:/Users/이윤정/Desktop/new.txt", 'w')
# f.write(read_ppt(prs))
# f.close()